from langchain_community.document_loaders import PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain.chains import RetrievalQA
from langchain_anthropic import ChatAnthropic
from langchain.schema import Document
from dotenv import load_dotenv
import os
import sys
import base64
import pandas as pd
from docx import Document as DocxDocument
from pptx import Presentation

sys.path.append(os.path.dirname(__file__))
load_dotenv()

# Global variables
vectorstore = None
rag_chain = None
uploaded_documents = []  # Track all uploaded documents

# Initialize embedding model once
embeddings = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")

# Initialize Claude once
llm = ChatAnthropic(
    model="claude-opus-4-6",
    anthropic_api_key=os.getenv("ANTHROPIC_API_KEY"),
    max_tokens=1024
)

def extract_text_from_file(file_path: str, filename: str) -> str:
    """Extract text from any supported file type"""
    ext = filename.lower().split(".")[-1]

    # PDF
    if ext == "pdf":
        from pypdf import PdfReader
        reader = PdfReader(file_path)
        return "\n".join([page.extract_text() or "" for page in reader.pages])

    # Word
    elif ext == "docx":
        doc = DocxDocument(file_path)
        return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])

    # Excel
    elif ext in ["xlsx", "xls"]:
        df_dict = pd.read_excel(file_path, sheet_name=None)
        text = ""
        for sheet_name, df in df_dict.items():
            text += f"\n\n--- Sheet: {sheet_name} ---\n"
            text += df.to_string(index=False)
        return text

    # PowerPoint
    elif ext == "pptx":
        prs = Presentation(file_path)
        text = ""
        for i, slide in enumerate(prs.slides):
            text += f"\n\n--- Slide {i+1} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
        return text

    # CSV
    elif ext == "csv":
        df = pd.read_csv(file_path)
        return df.to_string(index=False)

    # Plain text or markdown
    elif ext in ["txt", "md"]:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    # Image — Claude Vision OCR
    elif ext in ["jpg", "jpeg", "png", "webp"]:
        import anthropic
        client = anthropic.Anthropic()
        with open(file_path, "rb") as f:
            image_data = base64.b64encode(f.read()).decode()
        media_type = "image/jpeg" if ext in ["jpg", "jpeg"] else f"image/{ext}"
        response = client.messages.create(
            model="claude-opus-4-6",
            max_tokens=4096,
            messages=[{
                "role": "user",
                "content": [
                    {
                        "type": "image",
                        "source": {
                            "type": "base64",
                            "media_type": media_type,
                            "data": image_data
                        }
                    },
                    {
                        "type": "text",
                        "text": "Extract ALL content from this document image. Include all text, tables, charts descriptions, and any other visible information. Be thorough and preserve structure."
                    }
                ]
            }]
        )
        return response.content[0].text

    else:
        return f"Unsupported file type: {ext}"


def ingest_document(file_path: str, filename: str):
    """Read any document, chunk it, embed it, ADD to FAISS"""
    global vectorstore, rag_chain, uploaded_documents

    text = extract_text_from_file(file_path, filename)

    if not text.strip():
        return 0

    documents = [Document(
        page_content=text,
        metadata={"source": filename}
    )]

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=500,
        chunk_overlap=50
    )
    chunks = splitter.split_documents(documents)

    # ADD to existing FAISS or create new one
    if vectorstore is None:
        vectorstore = FAISS.from_documents(chunks, embeddings)
    else:
        vectorstore.add_documents(chunks)

    # Track uploaded documents
    if filename not in uploaded_documents:
        uploaded_documents.append(filename)

    retriever = vectorstore.as_retriever(search_kwargs={"k": 3})
    rag_chain = RetrievalQA.from_chain_type(
        llm=llm,
        chain_type="stuff",
        retriever=retriever,
        return_source_documents=True
    )

    return len(chunks)


def get_uploaded_documents():
    return uploaded_documents


def clear_documents():
    global vectorstore, rag_chain, uploaded_documents
    vectorstore = None
    rag_chain = None
    uploaded_documents = []


def ask_question(question: str):
    """Ask a question — FAISS first, Tavily fallback"""
    global rag_chain, vectorstore

    if rag_chain is None:
        return {"error": "No document uploaded yet. Please upload a document first."}

    # Step 1 — Search FAISS first
    result = rag_chain.invoke({"query": question})
    answer = result["result"]
    sources = len(result["source_documents"])

    # Step 2 — Check if FAISS found relevant answer
    no_answer_phrases = [
        "i don't have", "i do not have", "not mentioned",
        "no information", "not found", "cannot find",
        "not provided", "not available", "don't know"
    ]

    faiss_failed = any(phrase in answer.lower() for phrase in no_answer_phrases)

    # Step 3 — Tavily fallback if FAISS failed
    if faiss_failed:
        try:
            from tavily import TavilyClient
            tavily = TavilyClient(api_key=os.getenv("TAVILY_API_KEY"))
            search_results = tavily.search(query=question, max_results=3)

            web_context = "\n\n".join([
                f"Source: {r['url']}\n{r['content']}"
                for r in search_results.get("results", [])
            ])

            import anthropic
            client = anthropic.Anthropic()
            web_answer = client.messages.create(
                model="claude-opus-4-6",
                max_tokens=1024,
                messages=[{
                    "role": "user",
                    "content": f"Answer this question using the web sources below:\n\n{web_context}\n\nQuestion: {question}"
                }]
            )

            return {
                "answer": web_answer.content[0].text,
                "sources": len(search_results.get("results", [])),
                "source_type": "web"
            }
        except Exception as e:
            pass

    source_docs = list(set([
        doc.metadata.get("source", "Unknown")
        for doc in result["source_documents"]
    ]))

    return {
        "answer": answer,
        "sources": sources,
        "source_type": "document",
        "source_documents": source_docs
    }